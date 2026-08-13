"""Geração de apresentações e planilhas.

Divisão de trabalho que vale explicar: o modelo escreve o CONTEÚDO (título,
tópicos, dados) e este módulo faz a RENDERIZAÇÃO, de forma totalmente
determinística com python-pptx e openpyxl. Nenhuma chamada de modelo acontece
aqui.

Isso importa porque o estudo do Brahma-AI-Lite mostrou o caminho oposto sendo
tentado: lá existe um fluxo que raspa o Bing e um site de templates para baixar
arquivos .pptx prontos e preenchê-los. É frágil (raspagem quebra a cada mudança
de layout do site), depende de rede para criar um arquivo local, e redistribui
material de terceiros. Montar o tema em código não tem nenhum desses problemas
e o resultado fica consistente.

Os arquivos são criados apenas dentro das pastas permitidas e nunca sobrescrevem
nada — mesma regra do módulo de arquivos.
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from james.logs import audit, get_logger
from james.permissions.paths import PathNotAllowed
from james.tools.files import unique_destination
from james.tools.registry import Tool, ToolRegistry, ToolResult

logger = get_logger("james.tools.office")

MAX_SLIDES = 40
MAX_BULLETS_PER_SLIDE = 8
MAX_ROWS = 5000
MAX_COLUMNS = 60

# Paleta escura, na mesma linguagem visual do resto do James.
_BG = (0x07, 0x0D, 0x14)
_ACCENT = (0x00, 0xAA, 0xE1)
_TITLE = (0xBE, 0xF0, 0xFF)
_BODY = (0xC8, 0xDB, 0xE6)

_CHART_TYPES = {"barra": "bar", "coluna": "bar", "linha": "line", "pizza": "pie"}


def _slugify(text: str, default: str) -> str:
    """Nome de arquivo seguro a partir de um título falado."""
    cleaned = re.sub(r"[^\w\s-]", "", str(text or ""), flags=re.UNICODE).strip()
    cleaned = re.sub(r"[\s_-]+", "-", cleaned)
    return cleaned[:60] or default


def _output_path(guard, config, titulo: str, extensao: str) -> Path:
    """Resolve onde gravar, sempre dentro da whitelist.

    Sem uma pasta configurada, cair no diretório corrente seria um jeito
    silencioso de escrever fora da whitelist — melhor recusar e explicar.
    """
    configurada = config.resolve_path("office.output_dir")
    if configurada is None:
        permitidas = config.get("permissions.folders", []) or []
        if not permitidas:
            raise PathNotAllowed(
                "Nenhuma pasta de saída configurada. Defina office.output_dir "
                "ou permissions.folders no config.yaml."
            )
        configurada = Path(str(permitidas[0]))

    pasta = guard.validate_path(str(configurada))
    pasta.mkdir(parents=True, exist_ok=True)

    carimbo = datetime.now().strftime("%Y%m%d-%H%M")
    nome = f"{_slugify(titulo, 'documento')}-{carimbo}{extensao}"
    destino = unique_destination(pasta / nome)
    # Revalida o caminho final: o nome vem de texto do modelo.
    return guard.validate_path(str(destino))


def _normalize_slides(raw) -> list[dict]:
    """Aceita o formato do modelo e devolve algo previsível.

    Modelos entregam isto de formas variadas — lista de strings, lista de
    dicionários, chaves em português ou inglês. Normalizar aqui evita espalhar
    condicionais pelo renderizador.
    """
    slides: list[dict] = []
    for item in (raw or [])[:MAX_SLIDES]:
        if isinstance(item, str):
            slides.append({"titulo": item.strip(), "topicos": []})
            continue
        if not isinstance(item, dict):
            continue

        titulo = str(item.get("titulo") or item.get("title") or "").strip()
        topicos_raw = item.get("topicos") or item.get("bullets") or item.get("pontos") or []
        if isinstance(topicos_raw, str):
            topicos_raw = [linha for linha in topicos_raw.splitlines() if linha.strip()]

        topicos = [
            re.sub(r"^[-*•]\s*", "", str(t)).strip()
            for t in list(topicos_raw)[:MAX_BULLETS_PER_SLIDE]
            if str(t).strip()
        ]
        if titulo or topicos:
            slides.append({"titulo": titulo or "Slide", "topicos": topicos})
    return slides


def register(registry: ToolRegistry, config, guard) -> None:
    def criar_apresentacao(args: dict) -> ToolResult:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.util import Emu, Inches, Pt
        except ImportError:
            return ToolResult.failure(
                "Preciso da biblioteca python-pptx para isso. Instale com pip install python-pptx."
            )

        titulo = str(args.get("titulo", "")).strip()
        if not titulo:
            return ToolResult.failure("Preciso de um título para a apresentação.")
        subtitulo = str(args.get("subtitulo", "")).strip()
        slides = _normalize_slides(args.get("slides"))
        if not slides:
            return ToolResult.failure("Preciso do conteúdo dos slides.")

        try:
            destino = _output_path(guard, config, titulo, ".pptx")
        except PathNotAllowed as exc:
            return ToolResult.failure(str(exc))

        prs = Presentation()
        largura, altura = prs.slide_width, prs.slide_height
        layout_branco = prs.slide_layouts[6]

        def fundo(slide):
            forma = slide.shapes.add_shape(1, 0, 0, largura, altura)  # 1 = retângulo
            forma.fill.solid()
            forma.fill.fore_color.rgb = RGBColor(*_BG)
            forma.line.fill.background()
            # Sem isto o retângulo fica na frente do texto adicionado depois.
            slide.shapes._spTree.remove(forma._element)
            slide.shapes._spTree.insert(2, forma._element)

        def barra_accent(slide, topo):
            barra = slide.shapes.add_shape(1, Inches(0.6), topo, Inches(1.1), Pt(4))
            barra.fill.solid()
            barra.fill.fore_color.rgb = RGBColor(*_ACCENT)
            barra.line.fill.background()

        def texto(slide, esquerda, topo, larg, alt, conteudo, tamanho, cor, negrito=False):
            caixa = slide.shapes.add_textbox(esquerda, topo, larg, alt)
            quadro = caixa.text_frame
            quadro.word_wrap = True
            paragrafo = quadro.paragraphs[0]
            paragrafo.text = conteudo
            fonte = paragrafo.runs[0].font
            fonte.size = Pt(tamanho)
            fonte.bold = negrito
            fonte.color.rgb = RGBColor(*cor)
            return quadro

        # --- slide de capa ---
        capa = prs.slides.add_slide(layout_branco)
        fundo(capa)
        barra_accent(capa, Emu(int(altura * 0.40)))
        texto(capa, Inches(0.6), Emu(int(altura * 0.44)), largura - Inches(1.2),
              Inches(1.4), titulo, 40, _TITLE, negrito=True)
        if subtitulo:
            texto(capa, Inches(0.6), Emu(int(altura * 0.63)), largura - Inches(1.2),
                  Inches(0.8), subtitulo, 18, _BODY)

        # --- slides de conteúdo ---
        for item in slides:
            slide = prs.slides.add_slide(layout_branco)
            fundo(slide)
            barra_accent(slide, Inches(0.55))
            texto(slide, Inches(0.6), Inches(0.75), largura - Inches(1.2),
                  Inches(1.0), item["titulo"], 28, _TITLE, negrito=True)

            if not item["topicos"]:
                continue
            quadro = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.9), largura - Inches(1.6), altura - Inches(2.6)
            ).text_frame
            quadro.word_wrap = True
            for indice, topico in enumerate(item["topicos"]):
                paragrafo = quadro.paragraphs[0] if indice == 0 else quadro.add_paragraph()
                paragrafo.text = f"•  {topico}"
                paragrafo.space_after = Pt(10)
                fonte = paragrafo.runs[0].font
                fonte.size = Pt(16)
                fonte.color.rgb = RGBColor(*_BODY)

        try:
            prs.save(str(destino))
        except OSError as exc:
            logger.error("Falha ao gravar %s: %s", destino, exc)
            return ToolResult.failure("Não consegui gravar a apresentação.")

        audit("criar_apresentacao", arquivo=str(destino), slides=len(slides))
        return ToolResult(
            ok=True,
            speech=f"Apresentação criada com {len(slides) + 1} slides, em {destino.name}.",
            data={"arquivo": str(destino), "slides": len(slides) + 1},
        )

    def criar_planilha(args: dict) -> ToolResult:
        try:
            from openpyxl import Workbook
            from openpyxl.chart import BarChart, LineChart, PieChart, Reference
            from openpyxl.styles import Alignment, Font, PatternFill
            from openpyxl.utils import get_column_letter
        except ImportError:
            return ToolResult.failure(
                "Preciso da biblioteca openpyxl para isso. Instale com pip install openpyxl."
            )

        titulo = str(args.get("titulo", "")).strip()
        if not titulo:
            return ToolResult.failure("Preciso de um título para a planilha.")

        colunas = [str(c).strip() for c in (args.get("colunas") or []) if str(c).strip()]
        if not colunas:
            return ToolResult.failure("Preciso dos nomes das colunas.")
        colunas = colunas[:MAX_COLUMNS]

        linhas_raw = args.get("linhas") or []
        if not isinstance(linhas_raw, list):
            return ToolResult.failure("As linhas precisam vir como uma lista.")

        linhas = []
        for linha in linhas_raw[:MAX_ROWS]:
            if isinstance(linha, dict):
                # O modelo às vezes manda registros em vez de listas.
                linhas.append([linha.get(coluna, "") for coluna in colunas])
            elif isinstance(linha, (list, tuple)):
                linhas.append(list(linha)[: len(colunas)])
            else:
                linhas.append([linha])

        try:
            destino = _output_path(guard, config, titulo, ".xlsx")
        except PathNotAllowed as exc:
            return ToolResult.failure(str(exc))

        wb = Workbook()
        ws = wb.active
        ws.title = _slugify(titulo, "Dados")[:31].replace("-", " ") or "Dados"

        cabecalho_fundo = PatternFill("solid", fgColor="0D2231")
        cabecalho_fonte = Font(bold=True, color="BEF0FF")
        for indice, nome in enumerate(colunas, start=1):
            celula = ws.cell(row=1, column=indice, value=nome)
            celula.fill = cabecalho_fundo
            celula.font = cabecalho_fonte
            celula.alignment = Alignment(horizontal="center", vertical="center")

        for numero_linha, valores in enumerate(linhas, start=2):
            for numero_coluna, valor in enumerate(valores, start=1):
                ws.cell(row=numero_linha, column=numero_coluna, value=_coerce(valor))

        # Largura pelo conteúdo mais longo, com teto: uma célula com um texto
        # gigante deixaria a coluna inutilizável.
        for indice, nome in enumerate(colunas, start=1):
            maior = len(str(nome))
            for valores in linhas:
                if indice <= len(valores):
                    maior = max(maior, len(str(valores[indice - 1])))
            ws.column_dimensions[get_column_letter(indice)].width = min(45, maior + 4)

        ws.freeze_panes = "A2"

        grafico_spec = args.get("grafico") or {}
        if isinstance(grafico_spec, dict) and grafico_spec and len(linhas) > 0:
            tipo = _CHART_TYPES.get(
                str(grafico_spec.get("tipo", "barra")).strip().lower(), "bar"
            )
            classe = {"bar": BarChart, "line": LineChart, "pie": PieChart}[tipo]
            grafico = classe()
            grafico.title = str(grafico_spec.get("titulo") or titulo)
            # Primeira coluna vira rótulo, as demais viram séries.
            dados = Reference(
                ws, min_col=2, max_col=len(colunas), min_row=1, max_row=len(linhas) + 1
            )
            rotulos = Reference(ws, min_col=1, min_row=2, max_row=len(linhas) + 1)
            grafico.add_data(dados, titles_from_data=True)
            grafico.set_categories(rotulos)
            ws.add_chart(grafico, f"{get_column_letter(len(colunas) + 2)}2")

        try:
            wb.save(str(destino))
        except OSError as exc:
            logger.error("Falha ao gravar %s: %s", destino, exc)
            return ToolResult.failure("Não consegui gravar a planilha.")

        audit("criar_planilha", arquivo=str(destino), linhas=len(linhas))
        return ToolResult(
            ok=True,
            speech=f"Planilha criada com {len(linhas)} linhas, em {destino.name}.",
            data={"arquivo": str(destino), "linhas": len(linhas)},
        )

    registry.register(
        Tool(
            name="criar_apresentacao",
            description=(
                "Cria um arquivo PowerPoint. Você escreve o conteúdo: título, "
                "subtítulo e a lista de slides, cada um com título e tópicos. "
                "Use quando o usuário pedir uma apresentação ou slides."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "titulo": {"type": "string", "description": "Título da apresentação."},
                    "subtitulo": {"type": "string", "description": "Subtítulo da capa."},
                    "slides": {
                        "type": "array",
                        "description": "Slides de conteúdo, em ordem.",
                        "items": {
                            "type": "object",
                            "properties": {
                                "titulo": {"type": "string"},
                                "topicos": {"type": "array", "items": {"type": "string"}},
                            },
                            "required": ["titulo"],
                        },
                    },
                },
                "required": ["titulo", "slides"],
            },
            handler=criar_apresentacao,
            fire_and_forget=True,
        )
    )
    registry.register(
        Tool(
            name="criar_planilha",
            description=(
                "Cria um arquivo Excel com cabeçalho formatado e, opcionalmente, "
                "um gráfico. Use quando o usuário pedir uma planilha ou tabela."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "titulo": {"type": "string", "description": "Título da planilha."},
                    "colunas": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Nomes das colunas.",
                    },
                    "linhas": {
                        "type": "array",
                        "description": "Linhas de dados, cada uma uma lista de valores.",
                        "items": {"type": "array", "items": {"type": "string"}},
                    },
                    "grafico": {
                        "type": "object",
                        "description": "Gráfico opcional.",
                        "properties": {
                            "tipo": {"type": "string", "enum": ["barra", "linha", "pizza"]},
                            "titulo": {"type": "string"},
                        },
                    },
                },
                "required": ["titulo", "colunas", "linhas"],
            },
            handler=criar_planilha,
            fire_and_forget=True,
        )
    )


def _coerce(valor):
    """Converte texto numérico em número, para o Excel poder somar e plotar."""
    if isinstance(valor, (int, float, bool)) or valor is None:
        return valor
    texto = str(valor).strip()
    if not texto:
        return ""
    # Formato brasileiro: 1.234,56 -> 1234.56
    candidato = texto.replace("R$", "").replace("%", "").strip()
    if re.fullmatch(r"-?\d{1,3}(\.\d{3})*(,\d+)?", candidato):
        candidato = candidato.replace(".", "").replace(",", ".")
    try:
        numero = float(candidato)
    except ValueError:
        return texto
    return int(numero) if numero.is_integer() and "." not in candidato else numero
